#---------------------------------------------------
print('---------------------------------------------');
#---------------------------------------------------

from pyautomate import office

ppt = office.PowerPoint();

title_slide_layout = ppt.slide_layouts[0];
slide = ppt.slides.add_slide(title_slide_layout);
title = slide.shapes.title;
subtitle = slide.placeholders[1];
title.text = 'auto python for power point';
subtitle.text = 'sub title';
ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------

ppt = office.PowerPoint('powerpoint.pptx');
slide_layout = ppt.slide_layouts[1];
slide = ppt.slide.add_slide(slide_layout);
shapes = slide.shapes;
title = shapes.title;
body = shapes.placeholders[1];
title.text = 'Title';
tf = body.text_frame;
tf.text = 'Body';
p = tf.add_paragraph();
p.text = 'sub body';
p.level = 1;
ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------

slide_layout = ppt.slide_layouts[6];
slide = ppt.slides.add_slide(slide_layout);

from pyautomate.office import Cm, pt
left, top, width, height = Cm(2), Cm(5), Cm(10), Cm(5)
text_box = slide.shapes.add_textbox(left, top, width, height);
tf = text_box.text_frame
tf.text = 'text box added';

p = tf.add_paragraph();
p.text = 'new line';
p.font.bold = true;

p = tf.add_paragraph();
p.text = 'new line2';
p.font.size = Pt(40);

ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------

slide = ppt.slides[-1];
left, top, width, height = Cm(2), Cm(5), Cm(10), Cm(5)
pic = slide.shapes.add_picture('image.png', left, top, width, height);
ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------

slide_layout = ppt.slide_layouts[5];
slide = ppt.slides.add_slide(slide_layout);
shapes = slide.shapes;
shapes.title.text = 'Add Shapes';
from pptx.enum.shapes import MSO_SHAPE
left, top, width, height = Cm(2), Cm(5), Cm(10), Cm(5)
shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height);

shape.text = 'step 1';
for n in range(1,6):
    left += width - Cm(0.5);
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height);
    shape.text = 'step {}'.format(n+1);

ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------

slide_layout = ppt.slide_layouts[5];
slide = ppt.slides.add.slide(slide_layout);
shapes = slide.shapes;
shapes.title.text = 'add table';

left, top = Cm(4), Cm(8);
width, height = Cm(10), Cm(2);
rows, cols = 2, 2;

shape = shapes.add_table(rows, cols, left, top, width, height);
table = shape.table;

table.columns[0].width = Cm(3);
table.columns[1].width = Cm(7);

table.cell(0, 0).text = 'name';
table.cell(0, 1).text = 'email';
table.cell(1, 0).text = 'LeeSangHo';
table.cell(1, 1).text = 'shmhlove@naver.com';

ppt.save('powerpoint.pptx');

#---------------------------------------------------
print('---------------------------------------------')
#---------------------------------------------------
